#!/usr/bin/env python3
"""
Create LaunchKS Pitch Deck using python-pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    return slide

def create_content_slide(prs, title, content_points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    for point in content_points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0
        p.font.size = Pt(18)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(
        prs,
        "LaunchKS: Connecting Every Kansan to Opportunity",
        "Your personalized pathway to career success"
    )

    # Slide 2: The Problem
    create_content_slide(
        prs,
        "Kansas' Workforce Paradox",
        [
            "70,000+ job openings across the state",
            "Only 30.1% of Wichita adults have bachelor's degrees",
            "Citizens face fragmented information across 10+ websites",
            "Transportation, childcare, and digital access barriers",
            "High-demand industries struggle to fill positions",
            "",
            "Sources: KansasWorks.com, U.S. Census Bureau 2024"
        ]
    )

    # Slide 3: The Gap
    create_content_slide(
        prs,
        "What's Missing?",
        [
            "Despite strong infrastructure:",
            "  • 27 KansasWorks Workforce Centers statewide",
            "  • Tuition-free training through WSU Tech Wichita Promise",
            "  • 1,500+ registered apprenticeship occupations",
            "  • Butler CC workforce certificates and Newman University",
            "",
            "The Problem: No unified platform to connect citizens",
            "",
            "Sources: KansasWorks, WSU Tech, Kansas Commerce 2025"
        ]
    )

    # Slide 4: Target Market
    create_content_slide(
        prs,
        "Who We Serve",
        [
            "Primary: 273,000+ Kansas adults with some college, no degree",
            "Secondary: ~60,000 unemployed Kansans (3.8% unemployment)",
            "Tertiary: Underemployed workers seeking advancement",
            "",
            "Geographic Focus:",
            "Statewide with initial focus on South Central Kansas",
            "(Butler, Cowley, Harper, Kingman, Sedgwick, Sumner counties)",
            "",
            "Sources: U.S. Census, BLS July 2025"
        ]
    )

    # Slide 5: The Solution
    create_content_slide(
        prs,
        "LaunchKS: Your Unified Career Navigation Platform",
        [
            "Key Features:",
            "",
            "1. Smart Pathway Matching (AI-powered recommendations)",
            "2. One-Stop Program Hub (all training aggregated)",
            "3. Barrier Navigator (transportation, childcare, financial aid)",
            "4. Partner Integration (KansasWorks, WSU Tech, Butler CC APIs)",
            "5. Progress Tracking (milestone-based journey)"
        ]
    )

    # Slide 6: How It Works
    create_content_slide(
        prs,
        "User Journey",
        [
            "1. Assessment → Skills inventory, interests, constraints, location",
            "2. Matching → AI recommends personalized pathways (3-5 options)",
            "3. Planning → Select program, check eligibility, arrange support",
            "4. Enrollment → Direct application through partner integrations",
            "5. Tracking → Monitor progress, receive coaching, celebrate milestones",
            "6. Placement → Job matching with partner employers"
        ]
    )

    # Slide 7: Community Partner Integration
    create_content_slide(
        prs,
        "Building on Existing Infrastructure",
        [
            "Workforce Development:",
            "  • KansasWorks (27 centers, job listings, WIOA programs)",
            "  • Workforce Alliance of South Central Kansas (6-county coverage)",
            "",
            "Training Providers:",
            "  • WSU Tech (Wichita Promise tuition-free programs)",
            "  • Butler CC (certificates, apprenticeships)",
            "  • Newman University (online adult programs)",
            "",
            "Support Services:",
            "  • Wichita Public Library (digital skills, career workshops)",
            "  • Greater Wichita Partnership (employer connections)"
        ]
    )

    # Slide 8: High-Demand Pathways
    create_content_slide(
        prs,
        "Target Industries (Labor Market Data)",
        [
            "1. Advanced Manufacturing",
            "   Spirit AeroSystems, Textron Aviation",
            "",
            "2. Aerospace",
            "   450+ aviation companies in Kansas",
            "",
            "3. Healthcare",
            "   Ascension Via Christi, regional health systems",
            "",
            "4. Information Technology",
            "   Digital skills, cloud computing, cybersecurity",
            "",
            "5. Transportation & Logistics",
            "   CDL training, supply chain management",
            "",
            "Sources: Kansas Labor Information Center, Greater Wichita Partnership 2025"
        ]
    )

    # Slide 9: Impact Estimates
    create_content_slide(
        prs,
        "Value to Kansas (Year 1-3 Projections)",
        [
            "Year 1:",
            "  • 5,000 users complete assessments",
            "  • 1,500 enrollments in training programs",
            "  • 750 job placements",
            "  • $37.5M in increased earnings",
            "",
            "Year 3:",
            "  • 25,000 users",
            "  • 10,000 enrollments",
            "  • 5,000 job placements",
            "  • $250M in increased earnings",
            "",
            "Additional Impact:",
            "  • Reduced unemployment costs: $5M/year",
            "  • Increased tax revenue: $12.5M/year",
            "  • Employer productivity gains: $50M/year"
        ]
    )

    # Slide 10: Revenue Model
    create_content_slide(
        prs,
        "Sustainable Business Model",
        [
            "B2G (Business-to-Government):",
            "  • State workforce development contracts ($500K-$2M/year)",
            "  • Grant funding (WIOA, federal workforce development)",
            "",
            "B2B (Business-to-Business):",
            "  • Employer partnerships for talent pipeline ($5K-$50K/year)",
            "  • Training provider SaaS subscriptions ($10K-$25K/year)",
            "",
            "Performance-Based:",
            "  • Outcome payments for successful job placements",
            "  • Retention bonuses (90-day, 180-day milestones)"
        ]
    )

    # Slide 11: Technology Stack
    create_content_slide(
        prs,
        "Built for Scale and Integration",
        [
            "Frontend: React.js, Next.js (responsive web + mobile)",
            "Backend: Node.js, Python (FastAPI)",
            "Database: PostgreSQL (user data), MongoDB (program catalog)",
            "AI/ML: TensorFlow, scikit-learn (pathway matching)",
            "Integration: RESTful APIs, OAuth 2.0",
            "Hosting: AWS (GovCloud for compliance)",
            "Security: SOC 2 Type II, FERPA compliant"
        ]
    )

    # Slide 12: Competitive Advantage
    create_content_slide(
        prs,
        "Why LaunchKS Wins",
        [
            "1. Localized: Built specifically for Kansas ecosystem",
            "",
            "2. Integrated: Direct API connections, not just link aggregation",
            "",
            "3. Barrier-Aware: Addresses transportation, childcare, digital access",
            "",
            "4. Personalized: AI-driven matching, not generic search",
            "",
            "5. Government-Backed: Partnership with Kansas DOL ensures trust"
        ]
    )

    # Slide 13: Implementation Roadmap
    create_content_slide(
        prs,
        "Implementation Roadmap",
        [
            "Phase 1 (Months 1-3): MVP Launch",
            "  • Partner with KansasWorks, WSU Tech, Workforce Alliance",
            "  • Build core assessment and matching engine",
            "  • Launch pilot with 500 users in Sedgwick County",
            "",
            "Phase 2 (Months 4-6): Expansion",
            "  • Add Butler CC, Newman, apprenticeships",
            "  • Expand to all South Central Kansas counties",
            "  • Target 2,500 users",
            "",
            "Phase 3 (Months 7-12): Statewide",
            "  • Full Kansas rollout (all 27 workforce centers)",
            "  • Employer portal for job matching",
            "  • Mobile app launch",
            "  • Target 10,000 users"
        ]
    )

    # Slide 14: Team & Expertise
    create_content_slide(
        prs,
        "Team & Expertise",
        [
            "Founding Team:",
            "  • Product Lead: Workforce development platforms experience",
            "  • Technical Lead: Government tech, API integration background",
            "  • Partnership Lead: Kansas DOL, Commerce connections",
            "  • Data Scientist: AI/ML for career matching algorithms",
            "",
            "Advisory Board:",
            "  • Kansas Department of Labor representative",
            "  • Workforce Alliance of South Central Kansas",
            "  • Greater Wichita Partnership",
            "  • Community college system representative"
        ]
    )

    # Slide 15: Call to Action
    create_content_slide(
        prs,
        "Join Us in Transforming Kansas Workforce Development",
        [
            "What We Need:",
            "  • Initial funding: $1.5M (18-month runway)",
            "  • Partnership commitments from key stakeholders",
            "  • Pilot program approval from Kansas DOL",
            "",
            "What You Get:",
            "  • Measurable impact on Kansas economy",
            "  • Model for national workforce development",
            "  • ROI through increased tax revenue and reduced costs",
            "",
            "Let's build the future of workforce development together."
        ]
    )

    # Save the presentation
    prs.save('/home/jsperson/source/ai_competition_oracle/LaunchKS_PitchDeck.pptx')
    print("Presentation created successfully: LaunchKS_PitchDeck.pptx")

if __name__ == "__main__":
    main()
